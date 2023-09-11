from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

# Function to set font style for a run
def set_font(run, font_name, font_size):
    r = run._r
    rPr = r.get_or_add_rPr()
    rPr.append(parse_xml(f'<a:latin typeface="{font_name}" {nsdecls("a")} />'))
    rPr.append(parse_xml(f'<a:cs typeface="{font_name}" {nsdecls("a")} />'))
    rPr.append(parse_xml(f'<a:sz val="{font_size}" {nsdecls("a")} />'))

# Function to add content to a slide from a text file
def add_content_to_slide(slide, text_file, font_file):
    # Read content from text file
    with open(text_file, 'r') as file:
        content = file.read()
    
    # Add a textbox to the slide
    left = top = Pt(1)
    width = height = Pt(500)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    
    # Set textbox properties
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    # Set font properties
    font_name = font_file.split('/')[-1].split('.')[0]  # Extract font name from file path
    font_size = Pt(18)
    
    # Set font style for each paragraph in the content
    for paragraph in content.split('\n'):
        p = text_frame.add_paragraph()
        p.text = paragraph
        
        for run in p.runs:
            set_font(run, font_name, font_size)
    
    # Set font style for the title placeholder
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            title = shape.text_frame
            for paragraph in title.paragraphs:
                for run in paragraph.runs:
                    set_font(run, font_name, font_size)

try:
    # Create a new PowerPoint presentation
    presentation = Presentation()

    # Add two slides
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])

    # Add content to Slide 1
    add_content_to_slide(slide1, "page 1 ")

    # Add content to Slide 2
    add_content_to_slide(slide2, 'page 2')

    # Save the presentation
    output_file = 'output.pptx'
    presentation.save(output_file)
    print(f"Presentation saved as {output_file}")

except Exception as e:
    print(f"Error occurred: {str(e)}")
